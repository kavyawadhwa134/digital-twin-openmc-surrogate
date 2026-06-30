"""Build the ICTP workshop deck for the Digital Twin surrogate project.

All numbers are taken verbatim from committed evaluation JSONs:
  models/assembly_groupconst_evaluation.json
  models/pincell_lhs700_rigorous_evaluation.json
  models/state_forecaster_metrics.json
  models/xs_torch_benchmark_1M.json
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
OUT = ROOT / "ICTP_DigitalTwin_Surrogate.pptx"

# ---- Palette: Ocean Gradient (nuclear / physics) --------------------------
DEEP   = RGBColor(0x06, 0x3A, 0x5E)   # deep navy blue
BLUE   = RGBColor(0x1C, 0x72, 0x93)   # teal-blue
TEAL   = RGBColor(0x2C, 0x7A, 0x7B)   # teal
MID    = RGBColor(0x16, 0x20, 0x3A)   # midnight (title bg)
INK    = RGBColor(0x20, 0x28, 0x33)   # body text
MUTE   = RGBColor(0x5b, 0x6b, 0x7b)   # muted caption
LIGHT  = RGBColor(0xEE, 0xF3, 0xF7)   # light card
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
AMBER  = RGBColor(0xC9, 0x7B, 0x1A)
RED    = RGBColor(0xB3, 0x26, 0x1E)
ICEBG  = RGBColor(0xF6, 0xF9, 0xFB)

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s


def txt(s, x, y, w, h, lines, size=16, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font=BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
        space_after=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if isinstance(ln, tuple):
            text, opts = ln
        else:
            text, opts = ln, {}
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.name = opts.get("font", font)
        f.color.rgb = opts.get("color", color)
    return tb


def card(s, x, y, w, h, fill=LIGHT, line=None, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    return shp


def img_fit(s, path, x, y, w, h, align="center", valign="middle"):
    """Place image inside box (x,y,w,h) in inches, preserving aspect ratio."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    if align == "center":
        nx = x + (w - nw) / 2
    elif align == "left":
        nx = x
    else:
        nx = x + (w - nw)
    if valign == "middle":
        ny = y + (h - nh) / 2
    elif valign == "top":
        ny = y
    else:
        ny = y + (h - nh)
    s.shapes.add_picture(str(path), Inches(nx), Inches(ny), Inches(nw), Inches(nh))


def stat(s, x, y, w, number, label, color=BLUE, numsize=40, labsize=12):
    txt(s, x, y, w, 0.8, number, size=numsize, color=color, bold=True,
        font=HEAD, align=PP_ALIGN.CENTER)
    txt(s, x, y + 0.78, w, 0.6, label, size=labsize, color=MUTE,
        align=PP_ALIGN.CENTER, line_spacing=0.95)


def title_block(s, kicker, title, color=DEEP):
    txt(s, 0.7, 0.45, 12, 0.4, kicker, size=14, color=BLUE, bold=True, font=BODY)
    txt(s, 0.7, 0.78, 12, 1.0, title, size=32, color=color, bold=True, font=HEAD)


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide(MID)
card(s, 0, 0, 13.333, 7.5, fill=MID, radius=False)
txt(s, 0.9, 1.5, 11.5, 0.5, "ICTP WORKSHOP  ·  GENERATION IV REACTORS",
    size=15, color=RGBColor(0x9C, 0xC3, 0xDC), bold=True, font=BODY)
txt(s, 0.9, 2.1, 11.5, 2.0,
    "Machine-Learned Surrogate Modelling for\nDigital-Twin Applications and Accelerated\nMonte Carlo Neutron Transport",
    size=38, color=WHITE, bold=True, font=HEAD, line_spacing=1.05)
# divider via whitespace; no accent line
txt(s, 0.9, 5.25, 11.5, 0.5, "Kavya Wadhwa", size=22, color=WHITE, bold=True, font=BODY)
txt(s, 0.9, 5.85, 11.5, 0.5,
    "OpenMC  ·  physics-informed ML  ·  leakage-free evaluation",
    size=15, color=RGBColor(0x9C, 0xC3, 0xDC), font=BODY)
# three small pill markers
for i, (lab, c) in enumerate([("Pin cell", BLUE), ("Assembly", TEAL),
                              ("Digital twin", RGBColor(0x4A,0x6F,0x8A))]):
    px = 0.9 + i * 2.45
    p = card(s, px, 6.55, 2.2, 0.55, fill=c)
    txt(s, px, 6.62, 2.2, 0.45, lab, size=13, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)

# ===========================================================================
# 2. MOTIVATION
# ===========================================================================
s = slide(ICEBG)
title_block(s, "WHY", "Monte Carlo is the gold standard — and the bottleneck")
txt(s, 0.7, 1.85, 7.0, 3.6, [
    ("Monte Carlo neutron transport (OpenMC) is the reference for reactor "
     "physics — but a single converged calculation costs seconds to hours.", {}),
    ("", {"size": 6}),
    ("A digital twin must run thousands of evaluations: design sweeps, "
     "real-time state estimation, uncertainty quantification, control.", {}),
    ("", {"size": 6}),
    ("Idea: train an ML surrogate on OpenMC outputs, then query it instead. "
     "Pay the Monte Carlo cost once; amortize it over every later query.", {}),
    ("", {"size": 6}),
    ("This talk: three components, each evaluated honestly against the "
     "Monte Carlo noise floor — no overclaiming.", {"bold": True, "color": DEEP}),
], size=17, line_spacing=1.12, space_after=8)

cx = 8.1
card(s, cx, 1.9, 4.5, 4.6, fill=WHITE, line=RGBColor(0xD5,0xE0,0xE8))
txt(s, cx, 2.15, 4.5, 0.5, "The amortization argument", size=16, color=DEEP,
    bold=True, align=PP_ALIGN.CENTER, font=HEAD)
stat(s, cx+0.1, 2.85, 4.3, "53,000×", "assembly surrogate vs OpenMC", color=TEAL, numsize=44)
stat(s, cx+0.1, 4.25, 2.1, "~0.09 ms", "per ML query", color=BLUE, numsize=26, labsize=11)
stat(s, cx+2.3, 4.25, 2.1, "4.7 s", "per OpenMC case", color=MUTE, numsize=26, labsize=11)
txt(s, cx+0.2, 5.55, 4.1, 0.8,
    "Break-even ≈ the training-set size. Beyond that, the surrogate is pure speed.",
    size=12, color=MUTE, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.0)

# ===========================================================================
# 3. ROADMAP / LADDER
# ===========================================================================
s = slide(WHITE)
title_block(s, "THE ARC OF THE TALK", "A surrogate ladder toward whole-core simulation")
img_fit(s, FIG/"ictp_surrogate_ladder.png", 0.6, 1.9, 12.1, 4.7)
txt(s, 0.7, 6.75, 12, 0.5,
    "Rungs 1–2 are built and validated. Rung 3 — a nodal core solver consuming "
    "these group constants — is the defensible next step.",
    size=13, color=MUTE, italic=True, align=PP_ALIGN.CENTER)

# ===========================================================================
# 4. METHODOLOGY = HONESTY
# ===========================================================================
s = slide(ICEBG)
title_block(s, "HOW WE EVALUATE", "Rigor is the headline: three honesty guarantees")
items = [
    ("Leakage-free protocol",
     "Model selection by k-fold cross-validation on the training pool only. "
     "The test set is touched exactly once, after the model is frozen. "
     "Reported as repeated-holdout mean ± std."),
    ("The Monte Carlo noise floor is the yardstick",
     "OpenMC labels carry statistical noise (≈135–236 pcm here). A surrogate "
     "cannot — and should not — beat that floor. We report error relative to it."),
    ("Report relative MAE, not R²",
     "R² saturates near 1.0 for any reasonable model on smooth physics. "
     "We lead with relative MAE in pcm, plus physics-consistency checks."),
]
y = 1.95
for i, (h, b) in enumerate(items):
    # number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.7), Inches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = [DEEP, BLUE, TEAL][i]
    badge.line.fill.background(); badge.shadow.inherit = False
    txt(s, 0.8, y+0.12, 0.7, 0.5, str(i+1), size=22, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, font=HEAD)
    txt(s, 1.75, y-0.05, 10.9, 0.5, h, size=20, color=DEEP, bold=True, font=HEAD)
    txt(s, 1.75, y+0.5, 10.9, 1.0, b, size=15, color=INK, line_spacing=1.08)
    y += 1.55

# ===========================================================================
# 5. RUNG 1 - PIN CELL setup
# ===========================================================================
s = slide(WHITE)
title_block(s, "RUNG 1  ·  PIN CELL", "Fast response surrogate for a single fuel pin")
img_fit(s, FIG/"openmc_pincell_geometry_cross_section.png", 0.6, 1.9, 5.2, 4.8, align="left")
xc = 6.2
txt(s, xc, 1.95, 6.4, 4.8, [
    ("Inputs (8 design parameters):", {"bold": True, "color": DEEP, "size": 17}),
    ("fuel temperature, enrichment, moderator density & temperature, fuel "
     "radius, pin pitch, cladding thickness, boron — plus engineered geometry "
     "features.", {"size": 15}),
    ("", {"size": 8}),
    ("Targets (8 reactor responses):", {"bold": True, "color": DEEP, "size": 17}),
    ("k_eff, fuel & moderator flux, fission rate, fuel/moderator/total capture "
     "rate, power-density proxy.", {"size": 15}),
    ("", {"size": 8}),
    ("700-case Latin-hypercube OpenMC sweep, 4000 particles/case. "
     "Best model selected per target: a quadratic for k_eff, "
     "Gaussian processes for the field quantities.", {"size": 15}),
], line_spacing=1.1, space_after=6)

# ===========================================================================
# 6. RUNG 1 results
# ===========================================================================
s = slide(ICEBG)
title_block(s, "RUNG 1  ·  RESULTS", "Pin-cell k_eff lands on the noise floor")
img_fit(s, FIG/"pincell_lhs700_rigorous_keff_parity.png", 0.5, 1.95, 5.0, 4.7, align="left")
# stat strip right
sx = 6.0
stat(s, sx, 2.1, 3.4, "236 pcm", "k_eff MAE  (= 236 pcm MC floor)", color=BLUE, numsize=40)
stat(s, sx+3.5, 2.1, 3.3, "100%", "within 1000 pcm", color=GREEN, numsize=40)
stat(s, sx, 3.7, 3.4, "0.14–0.39%", "flux / rate field relMAE", color=TEAL, numsize=30, labsize=12)
stat(s, sx+3.5, 3.7, 3.3, "161 pcm", "intrinsic surrogate RMSE", color=BLUE, numsize=30, labsize=12)
card(s, sx, 5.0, 6.8, 1.55, fill=WHITE, line=RGBColor(0xD5,0xE0,0xE8))
txt(s, sx+0.25, 5.18, 6.3, 1.3, [
    ("Physics checks pass:", {"bold": True, "color": DEEP, "size": 15}),
    ("k_eff ↓ with boron, ↓ with fuel temperature (Doppler), ↑ with "
     "enrichment — all monotonic, as required.", {"size": 14}),
], line_spacing=1.05)

# ===========================================================================
# 7. RUNG 2 - ASSEMBLY (BEST)
# ===========================================================================
s = slide(WHITE)
title_block(s, "RUNG 2  ·  ASSEMBLY   ★ BEST RESULT", "From a pin to a real lattice — the bridge to cores")
img_fit(s, FIG/"assembly_geometry.png", 0.6, 1.9, 4.8, 4.9, align="left")
xc = 5.9
txt(s, xc, 1.95, 6.7, 4.9, [
    ("A 7×7 heterogeneous PWR assembly", {"bold": True, "color": DEEP, "size": 18}),
    ("UO₂ pins + 5 guide-tube water holes, reflective boundaries, run in "
     "OpenMC.", {"size": 15}),
    ("", {"size": 8}),
    ("The surrogate predicts 2-group homogenized constants", {"bold": True, "color": DEEP, "size": 18}),
    ("D₁, D₂, Σa, νΣf, Σf, the full scatter matrix, χ — and k∞. These are "
     "exactly the group constants a nodal / SP3 core solver consumes.", {"size": 15}),
    ("", {"size": 8}),
    ("Why this matters", {"bold": True, "color": TEAL, "size": 18}),
    ("This is the standard two-step (lattice → core) route. Surrogating this "
     "map is the real path from a pin to a full-core digital twin.", {"size": 15}),
], line_spacing=1.08, space_after=5)

# ===========================================================================
# 8. RUNG 2 results
# ===========================================================================
s = slide(ICEBG)
title_block(s, "RUNG 2  ·  RESULTS", "k∞ within 17 pcm of the noise floor; 53,000× faster")
img_fit(s, FIG/"assembly_kinf_parity.png", 0.5, 2.0, 4.6, 4.6, align="left")
sx = 5.5
stat(s, sx, 2.05, 3.6, "152 pcm", "k∞ MAE  (floor 135 pcm)", color=TEAL, numsize=40)
stat(s, sx+3.6, 2.05, 3.5, "53,000×", "speedup vs OpenMC", color=BLUE, numsize=40)
stat(s, sx, 3.65, 3.6, "0.04–0.14%", "13 group constants relMAE", color=TEAL, numsize=28, labsize=12)
stat(s, sx+3.6, 3.65, 3.5, "375 / 125", "leakage-free train / test", color=BLUE, numsize=28, labsize=12)
card(s, sx, 4.95, 7.1, 1.6, fill=WHITE, line=RGBColor(0xD5,0xE0,0xE8))
txt(s, sx+0.25, 5.12, 6.6, 1.4, [
    ("100% physically consistent:", {"bold": True, "color": DEEP, "size": 15}),
    ("D₁>D₂, Σa₂>Σa₁, νΣf₂>νΣf₁, downscatter>upscatter, χ₁≈1 on every test "
     "case. Only the tiny up-scatter term Σs2→1 is harder (1.4%) — and it is "
     "noise-dominated, not model-limited.", {"size": 13.5}),
], line_spacing=1.05)

# ===========================================================================
# 9. METRIC STORY (R2 audit)
# ===========================================================================
s = slide(WHITE)
title_block(s, "METHODOLOGICAL MATURITY", "Why we refuse to headline with R² = 0.99")
img_fit(s, FIG/"r2_metric_audit.png", 0.6, 1.95, 12.1, 4.0)
txt(s, 0.7, 6.15, 12, 0.9, [
    ("A reviewer who knows reactor physics is not impressed by R²=0.99 — a "
     "plain linear fit already gives it. Relative MAE shows the surrogate "
     "genuinely beats linear by 2.2–3.7×. We adversarially audited this "
     "(permutation control R²=−0.55; survives on residuals after removing "
     "enrichment + boron trends).", {"size": 13, "italic": True, "color": MUTE}),
], line_spacing=1.05)

# ===========================================================================
# 10. NOISE FLOOR
# ===========================================================================
s = slide(ICEBG)
title_block(s, "THE CENTRAL CLAIM", "Both surrogates are as accurate as the data allows")
img_fit(s, FIG/"ictp_noise_floor_summary.png", 2.3, 1.95, 8.7, 4.4)
txt(s, 0.7, 6.55, 12, 0.6,
    "To go below the floor you add Monte Carlo particles, not model capacity. "
    "That is the honest ceiling — and we are sitting on it.",
    size=14, color=DEEP, bold=True, align=PP_ALIGN.CENTER, italic=True)

# ===========================================================================
# 11. DIGITAL TWIN
# ===========================================================================
s = slide(WHITE)
title_block(s, "COMPONENT 3", "Digital-twin state forecaster + anomaly detection")
img_fit(s, FIG/"state_forecaster_validation_summary.png", 0.55, 2.0, 7.2, 3.6, align="left")
sx = 8.0
stat(s, sx, 2.0, 4.6, "1.00 / 0.88 / 0.93", "held-out precision / recall / F1", color=BLUE, numsize=27, labsize=12)
stat(s, sx, 3.35, 2.25, "~0", "false-positive rate", color=GREEN, numsize=34, labsize=11)
stat(s, sx+2.35, 3.35, 2.25, "2 s", "median alarm delay", color=TEAL, numsize=34, labsize=11)
stat(s, sx, 4.75, 4.6, "≈1 K / 17 pcm", "fuel-temp / k_eff forecast RMSE", color=BLUE, numsize=27, labsize=12)
txt(s, 0.6, 5.85, 12.1, 1.1, [
    ("MLP-ensemble forecaster trained on normal traces only; a hybrid detector "
     "flags departures. Calibrated on coolant-loss, rod-withdrawal and "
     "flux-bias faults; tested on unseen fault families and weak transients.", {"size": 13.5}),
    ("Honest caveat: trajectories are physics-informed simulations "
     "(point-kinetics + thermal-hydraulics ODEs), not real plant data.",
     {"size": 13.5, "bold": True, "color": AMBER}),
], line_spacing=1.06, space_after=4)

# ===========================================================================
# 12. XS-GPU honest result
# ===========================================================================
s = slide(ICEBG)
title_block(s, "THE ABSTRACT'S STATED HEADLINE — REPORTED HONESTLY",
            "Microscopic XS GPU surrogate: a negative result, not buried")
card(s, 0.7, 2.0, 5.9, 4.6, fill=WHITE, line=RGBColor(0xE2,0xC9,0xC9))
txt(s, 0.95, 2.2, 5.4, 0.5, "What we found", size=18, color=RED, bold=True, font=HEAD)
txt(s, 0.95, 2.85, 5.4, 3.6, [
    ("• Dense NN inference does not beat OpenMC's vectorized HDF5 cross-section "
     "lookup on CPU (≈0.06× — i.e. slower).", {}),
    ("", {"size": 6}),
    ("• Resonance region is genuinely hard for a smooth network: ~15% median, "
     "~72% p95 relative error in the macroscopic accumulation.", {}),
    ("", {"size": 6}),
    ("• OpenMC's per-nuclide interpolation is a strong, cache-friendly "
     "baseline. We say so plainly.", {}),
], size=14.5, color=INK, line_spacing=1.08)

card(s, 6.9, 2.0, 5.7, 4.6, fill=WHITE, line=RGBColor(0xC9,0xDD,0xD0))
txt(s, 7.15, 2.2, 5.2, 0.5, "Why this is still defensible", size=18, color=GREEN, bold=True, font=HEAD)
txt(s, 7.15, 2.85, 5.2, 3.6, [
    ("• A negative result, scoped correctly, is science — not a failure.", {}),
    ("", {"size": 6}),
    ("• The regime where GPU inference *can* win is the memory-bound, "
     "random-energy, many-nuclide pattern (XSBench's real cost) — future A100 work.", {}),
    ("", {"size": 6}),
    ("• Even then, the honest design is a UQ-gated hybrid: use the surrogate "
     "where it is confident, fall back to lookup in resonances.", {}),
    ("", {"size": 6}),
    ("• The defensible contributions are Rungs 1–2, where we win cleanly.", {"bold": True, "color": DEEP}),
], size=14, color=INK, line_spacing=1.06)

# ===========================================================================
# 13. SCORECARD
# ===========================================================================
s = slide(WHITE)
title_block(s, "SUMMARY", "Honest scorecard — what is defensible today")
img_fit(s, FIG/"ictp_scorecard.png", 0.6, 1.95, 12.1, 5.0)

# ===========================================================================
# 14. ROADMAP / NEXT
# ===========================================================================
s = slide(ICEBG)
title_block(s, "WHAT'S NEXT", "The credible path to a full-core digital twin")
roadmap = [
    ("Burnup / depletion axis", DEEP,
     "Add fuel depletion to the assembly sweep — the #1 core state variable. "
     "Blocked only on an external decay-chain file (openmc.deplete is installed)."),
    ("Nodal core solver (Rung 3)", BLUE,
     "Feed the surrogate group constants into a nodal / SP3 diffusion solver "
     "for a fast full-core flux / power map."),
    ("Real operational data", TEAL,
     "Replace synthetic trajectories with measured plant data to harden the "
     "anomaly detector beyond simulation."),
    ("GPU XS in its real regime", RGBColor(0x4A,0x6F,0x8A),
     "Re-run the cross-section surrogate on A100 in the memory-bound "
     "many-nuclide pattern, as a UQ-gated hybrid."),
]
y = 2.0
for i, (h, c, b) in enumerate(roadmap):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    yy = 2.0 + row * 2.35
    card(s, x, yy, 5.7, 2.05, fill=WHITE, line=RGBColor(0xD5,0xE0,0xE8))
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.25), Inches(yy+0.25), Inches(0.55), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = c; badge.line.fill.background()
    badge.shadow.inherit = False
    txt(s, x+0.25, yy+0.34, 0.55, 0.4, str(i+1), size=18, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, font=HEAD)
    txt(s, x+1.0, yy+0.28, 4.5, 0.5, h, size=17, color=c, bold=True, font=HEAD)
    txt(s, x+0.3, yy+0.95, 5.1, 1.0, b, size=13.5, color=INK, line_spacing=1.05)

# ===========================================================================
# 15. CONCLUSIONS
# ===========================================================================
s = slide(MID)
card(s, 0, 0, 13.333, 7.5, fill=MID, radius=False)
txt(s, 0.9, 0.9, 11.5, 0.5, "TAKEAWAYS", size=15, color=RGBColor(0x9C,0xC3,0xDC), bold=True)
txt(s, 0.9, 1.4, 11.5, 0.9, "Three honest results, one clear best",
    size=34, color=WHITE, bold=True, font=HEAD)
points = [
    ("Assembly group-constant surrogate is the result.",
     "k∞ to 152 pcm — at the Monte Carlo noise floor — with a 53,000× speedup, "
     "leakage-free, producing exactly what a core solver consumes."),
    ("Both surrogates are noise-floor-limited.",
     "Accuracy is capped by the physics data, not the model. That is the "
     "strongest defensible statement we can make."),
    ("We report honestly.",
     "Relative MAE not R², the XS-GPU negative result stated openly, synthetic "
     "twin data flagged. Rigor is the contribution."),
]
y = 2.65
for h, b in points:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y+0.08), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(0x4F,0xB0,0xC6); dot.line.fill.background()
    dot.shadow.inherit = False
    txt(s, 1.4, y-0.05, 11.0, 0.5, h, size=20, color=WHITE, bold=True, font=HEAD)
    txt(s, 1.4, y+0.45, 11.0, 0.8, b, size=15, color=RGBColor(0xC9,0xDC,0xE8),
        line_spacing=1.05)
    y += 1.4
txt(s, 0.9, 6.95, 11.5, 0.4, "Kavya Wadhwa  ·  ICTP Workshop  ·  Thank you",
    size=13, color=RGBColor(0x9C,0xC3,0xDC))

# ===========================================================================
# 16. BACKUP - reviewer Q&A
# ===========================================================================
s = slide(WHITE)
title_block(s, "BACKUP", "Anticipated reviewer questions")
qa = [
    ("“Isn't R²=0.99 trivial here?”",
     "Yes — that's why we don't headline it. We report relMAE vs the noise "
     "floor and show the surrogate beats a linear fit 2.2–3.7×."),
    ("“How do you know it isn't overfit / leaked?”",
     "Test set touched once; permutation control R²=−0.55; train→test gap "
     "0.02%; learning curve converged by ~100 samples."),
    ("“Is the speedup fair?”",
     "Break-even ≈ training-set size is stated. The training Monte Carlo cost "
     "is counted; beyond break-even the surrogate is net-faster."),
    ("“Why does the XS-GPU claim fail?”",
     "OpenMC's vectorized lookup is cache-friendly and hard to beat on CPU. "
     "Reported as a scoped negative result; GPU's real regime is future work."),
    ("“Is the digital twin real?”",
     "No — physics-informed simulated trajectories (point-kinetics + TH ODEs). "
     "Stated explicitly; real-data hardening is on the roadmap."),
    ("“Can this actually reach a core?”",
     "Yes — the assembly outputs are standard 2-group constants for a nodal "
     "solver. That is Rung 3, plus a burnup axis."),
]
y = 1.95
for i, (q, a) in enumerate(qa):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 6.25
    yy = 1.95 + row * 1.65
    txt(s, x, yy, 5.9, 0.5, q, size=15, color=DEEP, bold=True, font=HEAD)
    txt(s, x, yy+0.5, 5.9, 1.0, a, size=13, color=INK, line_spacing=1.04)
    y += 1.65

prs.save(str(OUT))
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
